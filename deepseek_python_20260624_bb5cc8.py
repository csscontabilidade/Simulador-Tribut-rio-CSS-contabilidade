# simulador_css_final_completo.py
# CSS Contabilidade – Simulador Tributário 2026/2027
# Carlos Antonio Stevanelli dos Santos – CRC/RO 8319/O

import streamlit as st
import pandas as pd
import json
import unicodedata
import zipfile
import datetime
from io import BytesIO
from dataclasses import dataclass
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

st.set_page_config(page_title="CSS Contabilidade – Simulador Tributário", page_icon="⚖️", layout="wide")

st.markdown("""
<style>
    .stApp { background-color: #f2f2f2; }
    .stButton>button { background-color: #1a3c5e; color: white; font-weight: bold; }
    .st-bb { border-bottom-color: #c9a94e; }
    h1, h2, h3, h4 { color: #1a3c5e; }
    .stMetric { background-color: white; padding: 10px; border-radius: 5px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); }
</style>
""", unsafe_allow_html=True)

ESCRITORIO = {
    'nome': 'CSS CONTABILIDADE',
    'contador': 'CARLOS ANTONIO STEVANELLI DOS SANTOS',
    'crc': 'CRC/RO 8319/O',
    'endereco': 'Av. Flamboyant, 501-B',
    'telefone': '(69) 98491-1067',
    'email': 'csscontabilidade@email.com.br',
    'formacao': [
        'Bacharel em Ciências Contábeis – UNITINS',
        'Especialização em Direito Tributário – Faculdade Farol',
        'CPA 10 – AMBIMA',
        'Coach Integral Sistêmico',
        'SIS Assessment – Perfil Comportamental'
    ],
    'experiencia': [
        '2002-2004: Escritório de Contabilidade',
        '2004-2011: Instituição Financeira Cooperativa',
        '2011-2026: Escritório de Contabilidade',
        'Recuperação de Créditos Tributários',
        'Processos Administrativos SEFIN/RFB',
        'Processos Judiciais Trabalhistas'
    ]
}

ALIQUOTA_CBS = 0.088

TABELA_SIMPLES = {
    'I': [(180000,0.04,0), (360000,0.073,5940), (720000,0.095,13860), (1800000,0.107,22500), (3600000,0.143,87300), (4800000,0.19,378000)],
    'II': [(180000,0.045,0), (360000,0.09,8100), (720000,0.102,15840), (1800000,0.14,35640), (3600000,0.22,125640), (4800000,0.33,378000)],
    'III': [(180000,0.06,0), (360000,0.112,9360), (720000,0.135,17640), (1800000,0.16,35640), (3600000,0.21,125640), (4800000,0.33,648000)],
    'IV': [(180000,0.045,0), (360000,0.09,8100), (720000,0.102,15840), (1800000,0.14,35640), (3600000,0.22,125640), (4800000,0.33,378000)],
    'V': [(180000,0.155,0), (360000,0.18,4500), (720000,0.195,9900), (1800000,0.205,17100), (3600000,0.23,62100), (4800000,0.305,540000)]
}

PIS_COFINS_SIMPLES = {
    180000: 0.0275, 360000: 0.0295, 720000: 0.0315,
    1800000: 0.0330, 3600000: 0.0345, 4800000: 0.0365
}

def aliquota_efetiva_simples(rbt12, anexo):
    if rbt12 <= 0: return 0
    tab = TABELA_SIMPLES.get(anexo, TABELA_SIMPLES['I'])
    for faixa in tab:
        if rbt12 <= faixa[0]:
            return max(0, (rbt12 * faixa[1] - faixa[2]) / rbt12)
    return max(0, (rbt12 * tab[-1][1] - tab[-1][2]) / rbt12)

def fator_pis_cofins_simples(rbt12):
    limites = sorted(PIS_COFINS_SIMPLES.keys())
    for lim in limites:
        if rbt12 <= lim:
            return PIS_COFINS_SIMPLES[lim]
    return PIS_COFINS_SIMPLES[limites[-1]]

@dataclass
class Empresa:
    nome: str = ''
    anexo_simples: str = 'I'
    regime_atual: str = 'Simples Nacional'
    faturamento_mensal: float = 0.0
    folha_mensal: float = 0.0
    despesas_admin_mensal: float = 0.0
    outros_impostos_mensal: float = 0.0
    juros_mensal: float = 0.0
    compras_mensais: float = 0.0
    cmv_percentual: float = 0.0
    fornecedor_simples_pct: float = 30.0
    fornecedor_real_presumido_pct: float = 70.0
    perc_eletronico: float = 70.0
    perc_b2b: float = 20.0

    @property
    def fat_anual(self): return self.faturamento_mensal * 12
    @property
    def folha_anual(self): return self.folha_mensal * 13
    @property
    def compras_anuais(self): return self.compras_mensais * 12
    @property
    def rbt12(self): return self.fat_anual

class Simulador2026:
    def __init__(self, emp: Empresa):
        self.e = emp
    def _credito_pis_cofins(self):
        return self.e.compras_anuais * (self.e.fornecedor_real_presumido_pct/100) * 0.0925
    def simples(self):
        rbt = self.e.rbt12
        aliq = aliquota_efetiva_simples(rbt, self.e.anexo_simples)
        carga = rbt * aliq
        lucro = rbt - self.e.compras_anuais - self.e.folha_anual - self.e.despesas_admin_mensal*12 - self.e.outros_impostos_mensal*12 - self.e.juros_mensal*12 - carga
        return {'regime': 'Simples Nacional', 'carga_total': carga, 'aliquota_efetiva': aliq*100, 'lucro_apos': max(0, lucro), 'detalhes': {'Simples unificado': carga}}
    def presumido(self):
        rec = self.e.fat_anual
        base_irpj = rec * 0.08
        irpj = base_irpj * 0.15
        ad_irpj = max(0, (base_irpj - 240000) * 0.10)
        csll = rec * 0.12 * 0.09
        pis_cofins = rec * 0.0365
        enc_folha = self.e.folha_anual * 0.258
        carga = irpj + ad_irpj + csll + pis_cofins + enc_folha + self.e.outros_impostos_mensal*12
        lucro = rec - self.e.compras_anuais - self.e.folha_anual - enc_folha - self.e.despesas_admin_mensal*12 - self.e.juros_mensal*12 - irpj - ad_irpj - csll - pis_cofins - self.e.outros_impostos_mensal*12
        return {'regime': 'Lucro Presumido', 'carga_total': carga, 'aliquota_efetiva': (carga/rec)*100, 'lucro_apos': max(0, lucro), 'detalhes': {'IRPJ/CSLL': irpj+ad_irpj+csll, 'PIS/COFINS': pis_cofins, 'Encargos Folha': enc_folha}}
    def real(self):
        rec = self.e.fat_anual
        lucro_bruto = rec - self.e.compras_anuais
        desp = self.e.folha_anual + self.e.despesas_admin_mensal*12 + self.e.juros_mensal*12
        base_real = lucro_bruto - desp
        irpj = max(0, base_real * 0.15)
        ad_irpj = max(0, (base_real - 240000) * 0.10)
        csll = max(0, base_real * 0.09)
        deb_pis_cofins = rec * 0.0925
        cred_pis_cofins = self._credito_pis_cofins()
        pis_cofins = max(0, deb_pis_cofins - cred_pis_cofins)
        enc_folha = self.e.folha_anual * 0.258
        carga = irpj + ad_irpj + csll + pis_cofins + enc_folha + self.e.outros_impostos_mensal*12
        lucro = base_real - irpj - ad_irpj - csll - pis_cofins - self.e.outros_impostos_mensal*12
        return {'regime': 'Lucro Real', 'carga_total': carga, 'aliquota_efetiva': (carga/rec)*100, 'lucro_apos': max(0, lucro), 'detalhes': {'IRPJ/CSLL': irpj+ad_irpj+csll, 'PIS/COFINS líquido': pis_cofins, 'Encargos Folha': enc_folha}}

class Simulador2027:
    def __init__(self, emp: Empresa):
        self.e = emp
    def _credito_cbs(self):
        return self.e.compras_anuais * (self.e.fornecedor_real_presumido_pct/100) * ALIQUOTA_CBS
    def simples_hibrido(self):
        rec = self.e.fat_anual
        deb_cbs = rec * ALIQUOTA_CBS
        cred_cbs = self._credito_cbs()
        cbs_liquido = max(0, deb_cbs - cred_cbs)
        nova_aliq = max(0, aliquota_efetiva_simples(self.e.rbt12, self.e.anexo_simples) - fator_pis_cofins_simples(self.e.rbt12))
        simples_restante = rec * nova_aliq
        carga = simples_restante + cbs_liquido + self.e.outros_impostos_mensal*12
        lucro = rec - self.e.compras_anuais - self.e.folha_anual - self.e.despesas_admin_mensal*12 - self.e.juros_mensal*12 - cbs_liquido - simples_restante - self.e.outros_impostos_mensal*12
        return {'regime': 'Simples Híbrido (2027)', 'carga_total': carga, 'aliquota_efetiva': (carga/rec)*100, 'lucro_apos': max(0, lucro), 'detalhes': {'CBS (débito)': deb_cbs, 'Crédito CBS': cred_cbs, 'CBS líquido': cbs_liquido, 'Simples reduzido': simples_restante}}
    def presumido(self):
        rec = self.e.fat_anual
        base_irpj = rec * 0.08
        irpj = base_irpj * 0.15
        ad_irpj = max(0, (base_irpj - 240000) * 0.10)
        csll = rec * 0.12 * 0.09
        deb_cbs = rec * ALIQUOTA_CBS
        cred_cbs = self._credito_cbs()
        cbs_liquido = max(0, deb_cbs - cred_cbs)
        enc_folha = self.e.folha_anual * 0.258
        carga = irpj + ad_irpj + csll + cbs_liquido + enc_folha + self.e.outros_impostos_mensal*12
        lucro = rec - self.e.compras_anuais - self.e.folha_anual - enc_folha - self.e.despesas_admin_mensal*12 - self.e.juros_mensal*12 - irpj - ad_irpj - csll - cbs_liquido - self.e.outros_impostos_mensal*12
        return {'regime': 'Lucro Presumido (2027)', 'carga_total': carga, 'aliquota_efetiva': (carga/rec)*100, 'lucro_apos': max(0, lucro), 'detalhes': {'IRPJ/CSLL': irpj+ad_irpj+csll, 'CBS (débito)': deb_cbs, 'Crédito CBS': cred_cbs, 'CBS líquido': cbs_liquido, 'Encargos Folha': enc_folha}}
    def real(self):
        rec = self.e.fat_anual
        lucro_bruto = rec - self.e.compras_anuais
        desp = self.e.folha_anual + self.e.despesas_admin_mensal*12 + self.e.juros_mensal*12
        base_real = lucro_bruto - desp
        irpj = max(0, base_real * 0.15)
        ad_irpj = max(0, (base_real - 240000) * 0.10)
        csll = max(0, base_real * 0.09)
        deb_cbs = rec * ALIQUOTA_CBS
        cred_cbs = self._credito_cbs()
        cbs_liquido = max(0, deb_cbs - cred_cbs)
        enc_folha = self.e.folha_anual * 0.258
        carga = irpj + ad_irpj + csll + cbs_liquido + enc_folha + self.e.outros_impostos_mensal*12
        lucro = base_real - irpj - ad_irpj - csll - cbs_liquido - self.e.outros_impostos_mensal*12
        return {'regime': 'Lucro Real (2027)', 'carga_total': carga, 'aliquota_efetiva': (carga/rec)*100, 'lucro_apos': max(0, lucro), 'detalhes': {'IRPJ/CSLL': irpj+ad_irpj+csll, 'CBS (débito)': deb_cbs, 'Crédito CBS': cred_cbs, 'CBS líquido': cbs_liquido, 'Encargos Folha': enc_folha}}
    def split_payment(self):
        fat_mensal = self.e.faturamento_mensal
        vendas_eletr = fat_mensal * self.e.perc_eletronico / 100
        split = vendas_eletr * ALIQUOTA_CBS
        receb_liq = fat_mensal - split
        nec_bruta = split * 2
        credito_mensal = self._credito_cbs() / 12
        nec_liquida = max(0, nec_bruta - credito_mensal)
        return {'vendas_eletr': vendas_eletr, 'split': split, 'receb_liq': receb_liq, 'nec_bruta': nec_bruta, 'credito_mensal': credito_mensal, 'nec_liquida': nec_liquida}

class RelatorioPDF:
    def __init__(self, empresa, res26, res27, sp):
        self.emp = empresa
        self.res26 = res26
        self.res27 = res27
        self.sp = sp

    def gerar(self):
        def limpar(texto):
            if not isinstance(texto, str):
                texto = str(texto)
            nfkd = unicodedata.normalize('NFKD', texto)
            return nfkd.encode('ascii', 'ignore').decode('ascii')

        pdf = FPDF()
        pdf.add_page()
        pdf.set_fill_color(26, 60, 94)
        pdf.rect(0, 0, 210, 35, 'F')
        pdf.set_y(8)
        pdf.set_font("Arial", "B", 16)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 8, limpar(ESCRITORIO['nome']), ln=True, align="C")
        pdf.set_font("Arial", "", 9)
        pdf.cell(0, 6, limpar(f"Contador: {ESCRITORIO['contador']} | {ESCRITORIO['crc']}"), ln=True, align="C")
        pdf.set_y(40)
        pdf.set_text_color(0, 0, 0)

        pdf.set_font("Arial", "B", 12)
        pdf.cell(0, 8, limpar("1. DADOS DA EMPRESA"), ln=True)
        pdf.set_font("Arial", "", 10)
        info = [
            limpar(f"Nome: {self.emp.nome}"),
            limpar(f"Regime Atual: {self.emp.regime_atual}"),
            limpar(f"Faturamento Mensal: R$ {self.emp.faturamento_mensal:,.2f}   Anual: R$ {self.emp.fat_anual:,.2f}"),
            limpar(f"Folha Mensal: R$ {self.emp.folha_mensal:,.2f}   Anual (13º): R$ {self.emp.folha_anual:,.2f}")
        ]
        for i in info:
            pdf.cell(0, 6, i, ln=True)
        pdf.ln(4)

        pdf.set_font("Arial", "B", 12)
        pdf.cell(0, 8, limpar("2. CENÁRIO ATUAL (2026)"), ln=True)
        for r in self.res26.values():
            pdf.set_font("Arial", "B", 10)
            pdf.cell(0, 7, limpar(r['regime']), ln=True)
            pdf.set_font("Arial", "", 9)
            pdf.cell(0, 6, limpar(f"Carga Total: R$ {r['carga_total']:,.2f}   Alíquota efetiva: {r['aliquota_efetiva']:.2f}%"), ln=True)
            for k, v in r['detalhes'].items():
                pdf.cell(0, 6, limpar(f"  {k}: R$ {v:,.2f}"), ln=True)
            pdf.ln(2)

        pdf.set_font("Arial", "B", 12)
        pdf.cell(0, 8, limpar("3. CENÁRIO 2027 (COM CBS)"), ln=True)
        for r in self.res27.values():
            pdf.set_font("Arial", "B", 10)
            pdf.cell(0, 7, limpar(r['regime']), ln=True)
            pdf.set_font("Arial", "", 9)
            pdf.cell(0, 6, limpar(f"Carga Total: R$ {r['carga_total']:,.2f}   Alíquota efetiva: {r['aliquota_efetiva']:.2f}%"), ln=True)
            for k, v in r['detalhes'].items():
                pdf.cell(0, 6, limpar(f"  {k}: R$ {v:,.2f}"), ln=True)
            pdf.ln(2)

        pdf.set_font("Arial", "B", 12)
        pdf.cell(0, 8, limpar("4. SPLIT PAYMENT (CBS)"), ln=True)
        pdf.set_font("Arial", "", 10)
        pdf.cell(0, 6, limpar(f"CBS retida mensal: R$ {self.sp['split']:,.2f}"), ln=True)
        pdf.cell(0, 6, limpar(f"Recebimento líquido: R$ {self.sp['receb_liq']:,.2f}"), ln=True)
        pdf.cell(0, 6, limpar(f"Capital de giro sugerido: R$ {self.sp['nec_liquida']:,.2f}"), ln=True)

        pdf.add_page()
        pdf.set_font("Arial", "B", 12)
        pdf.cell(0, 8, limpar("5. CONSULTOR RESPONSÁVEL"), ln=True)
        pdf.set_font("Arial", "", 10)
        pdf.cell(0, 6, limpar(f"{ESCRITORIO['contador']} – {ESCRITORIO['crc']}"), ln=True)
        pdf.ln(3)
        pdf.set_font("Arial", "B", 10)
        pdf.cell(0, 6, limpar("Formação:"), ln=True)
        pdf.set_font("Arial", "", 9)
        for f in ESCRITORIO['formacao']:
            pdf.cell(0, 6, limpar(f"• {f}"), ln=True)
        pdf.ln(2)
        pdf.set_font("Arial", "B", 10)
        pdf.cell(0, 6, limpar("Experiência:"), ln=True)
        for e in ESCRITORIO['experiencia']:
            pdf.cell(0, 6, limpar(f"• {e}"), ln=True)

        return pdf.output(dest='S')

class SlidesPPT:
    def __init__(self, empresa, res26, res27, sp):
        self.emp = empresa
        self.res26 = res26
        self.res27 = res27
        self.sp = sp

    def gerar(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        AZUL = RGBColor(26, 60, 94)
        DOURADO = RGBColor(201, 169, 78)
        BRANCO = RGBColor(255, 255, 255)

        def add_slide(title):
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tb.text_frame.paragraphs[0].text = title
            tb.text_frame.paragraphs[0].font.size = Pt(24)
            tb.text_frame.paragraphs[0].font.bold = True
            tb.text_frame.paragraphs[0].font.color.rgb = AZUL
            return sl

        sl = prs.slides.add_slide(prs.slide_layouts[6])
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = AZUL
        tb = sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
        tf = tb.text_frame
        p = tf.paragraphs[0]; p.text = ESCRITORIO['nome']; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = DOURADO; p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph(); p.text = f"CRC {ESCRITORIO['crc']}"; p.font.size = Pt(14); p.font.color.rgb = BRANCO; p.alignment = PP_ALIGN.CENTER
        tb2 = sl.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(2))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]; p2.text = f"Simulação Tributária 2026/2027 – {self.emp.nome}"; p2.font.size = Pt(24); p2.font.color.rgb = DOURADO; p2.alignment = PP_ALIGN.CENTER

        sl = add_slide("CENÁRIO ATUAL (2026)")
        y = 2
        for r in self.res26.values():
            tb = sl.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.5))
            tb.text_frame.paragraphs[0].text = f"{r['regime']}: R$ {r['carga_total']:,.2f} ({r['aliquota_efetiva']:.1f}%)"
            tb.text_frame.paragraphs[0].font.size = Pt(14)
            y += 0.6

        sl = add_slide("CENÁRIO 2027 (COM CBS)")
        y = 2
        for r in self.res27.values():
            tb = sl.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.5))
            tb.text_frame.paragraphs[0].text = f"{r['regime']}: R$ {r['carga_total']:,.2f} ({r['aliquota_efetiva']:.1f}%)"
            tb.text_frame.paragraphs[0].font.size = Pt(14)
            y += 0.6

        sl = add_slide("SPLIT PAYMENT CBS")
        y = 2
        for k, v in self.sp.items():
            tb = sl.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.5))
            tb.text_frame.paragraphs[0].text = f"{k}: R$ {v:,.2f}"
            tb.text_frame.paragraphs[0].font.size = Pt(14)
            y += 0.6

        sl = add_slide("CONSULTOR RESPONSÁVEL")
        y = 2
        tb = sl.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.5))
        tb.text_frame.paragraphs[0].text = f"{ESCRITORIO['contador']} – {ESCRITORIO['crc']}"
        tb.text_frame.paragraphs[0].font.size = Pt(16)
        y += 1
        for f in ESCRITORIO['formacao']:
            tb = sl.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.4))
            tb.text_frame.paragraphs[0].text = f"• {f}"
            tb.text_frame.paragraphs[0].font.size = Pt(12)
            y += 0.4
        y += 0.5
        for e in ESCRITORIO['experiencia']:
            tb = sl.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.4))
            tb.text_frame.paragraphs[0].text = f"• {e}"
            tb.text_frame.paragraphs[0].font.size = Pt(12)
            y += 0.4

        out = BytesIO()
        prs.save(out)
        return out.getvalue()

def main():
    with st.sidebar:
        st.image("https://cdn-icons-png.flaticon.com/512/3135/3135715.png", width=80)
        st.markdown(f"### {ESCRITORIO['nome']}")
        st.markdown(f"*{ESCRITORIO['crc']}*")
        st.markdown(f"📞 {ESCRITORIO['telefone']}")

        with st.expander("📋 Dados da Empresa", expanded=True):
            nome = st.text_input("Nome", "Empresa Exemplo")
            regime_atual = st.selectbox("Regime Atual", ["Simples Nacional", "Lucro Presumido", "Lucro Real"])
            anexo = st.selectbox("Anexo Simples", ['I','II','III','IV','V']) if regime_atual == "Simples Nacional" else 'I'

        with st.expander("💰 Valores Mensais"):
            fat = st.number_input("Faturamento (R$)", 0.0, 10000000.0, 50000.0, 1000.0, format="%.2f")
            folha = st.number_input("Folha Pagamento (R$)", 0.0, 5000000.0, 8000.0, 500.0, format="%.2f")
            desp = st.number_input("Desp. Admin/Operac. (R$)", 0.0, 5000000.0, 3000.0, 500.0, format="%.2f")
            out_imp = st.number_input("Outros Impostos (R$)", 0.0, 5000000.0, 1000.0, 500.0, format="%.2f")
            juros = st.number_input("Juros (R$)", 0.0, 5000000.0, 500.0, 100.0, format="%.2f")
            compras = st.number_input("Compras (R$)", 0.0, 5000000.0, 25000.0, 1000.0, format="%.2f")
            cmv = st.slider("CMV %", 0, 100, 50)

        with st.expander("🏭 Fornecedores (% das compras)"):
            f_simples = st.slider("Fornecedor Simples Nacional (sem crédito)", 0, 100, 30)
            f_real_pres = st.slider("Fornecedor Lucro Real/Presumido (com crédito)", 0, 100, 70)

        perc_eletr = st.slider("Vendas Eletrônicas %", 0, 100, 70)
        perc_b2b = st.slider("Vendas B2B %", 0, 100, 20)

        st.caption(f"📅 Anual: Faturamento R$ {fat*12:,.2f} | Folha (13º) R$ {folha*13:,.2f}")

        if st.button("📤 Exportar Questionário"):
            dados = {'nome': nome, 'anexo_simples': anexo, 'regime_atual': regime_atual,
                     'faturamento_mensal': fat, 'folha_mensal': folha, 'despesas_admin_mensal': desp,
                     'outros_impostos_mensal': out_imp, 'juros_mensal': juros, 'compras_mensais': compras,
                     'cmv_percentual': cmv, 'fornecedor_simples_pct': f_simples,
                     'fornecedor_real_presumido_pct': f_real_pres, 'perc_eletronico': perc_eletr, 'perc_b2b': perc_b2b}
            st.download_button("⬇️ Baixar Questionário JSON", data=json.dumps(dados, indent=4, ensure_ascii=False), file_name="questionario.json", mime="application/json")

        uploaded_file = st.file_uploader("📥 Importar Questionário", type=["json"])
        if uploaded_file is not None:
            try:
                dados_import = json.load(uploaded_file)
                for key, value in dados_import.items():
                    st.session_state[key] = value
                st.success("Questionário carregado! Recarregue a página para aplicar os dados (funcionalidade em melhoria).")
            except:
                st.error("Arquivo inválido.")

        if st.button("📥 Baixar Aplicativo (Instalador)"):
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, "w") as zf:
                zf.writestr("simulador.py", open(__file__, "r", encoding="utf-8").read())
                zf.writestr("requirements.txt", "streamlit\npandas\nplotly\nfpdf\npython-pptx")
                zf.writestr("instalar.bat", "@echo off\npip install -r requirements.txt\nstart streamlit run simulador.py")
                zf.writestr("instalar.sh", "#!/bin/bash\npip install -r requirements.txt\nstreamlit run simulador.py")
            st.download_button("⬇️ Baixar ZIP com instalador", data=zip_buffer.getvalue(), file_name="CSS_Simulador.zip")

    empresa = Empresa(
        nome=nome, anexo_simples=anexo, regime_atual=regime_atual,
        faturamento_mensal=fat, folha_mensal=folha, despesas_admin_mensal=desp,
        outros_impostos_mensal=out_imp, juros_mensal=juros, compras_mensais=compras,
        cmv_percentual=cmv, fornecedor_simples_pct=f_simples,
        fornecedor_real_presumido_pct=f_real_pres, perc_eletronico=perc_eletr, perc_b2b=perc_b2b
    )

    tab1, tab2 = st.tabs(["📊 Simulação 2026 vs 2027", "📄 Relatórios"])

    with tab1:
        if st.button("Executar Simulação Completa"):
            sim26 = Simulador2026(empresa)
            sim27 = Simulador2027(empresa)
            res26 = {'Simples': sim26.simples(), 'Presumido': sim26.presumido(), 'Real': sim26.real()}
            res27 = {'Simples Híbrido': sim27.simples_hibrido(), 'Presumido': sim27.presumido(), 'Real': sim27.real()}
            sp = sim27.split_payment()

            st.subheader("🔹 2026 (atual)")
            df26 = pd.DataFrame(res26.values())
            st.dataframe(df26[['regime','carga_total','aliquota_efetiva','lucro_apos']].style.format({
                'carga_total': 'R$ {:,.2f}', 'aliquota_efetiva': '{:.2f}%', 'lucro_apos': 'R$ {:,.2f}'
            }))
            with st.expander("Detalhes 2026"):
                for r in res26.values():
                    st.write(f"**{r['regime']}**")
                    for k, v in r['detalhes'].items():
                        st.write(f"- {k}: R$ {v:,.2f}")

            st.subheader("🔸 2027 (com CBS)")
            df27 = pd.DataFrame(res27.values())
            st.dataframe(df27[['regime','carga_total','aliquota_efetiva','lucro_apos']].style.format({
                'carga_total': 'R$ {:,.2f}', 'aliquota_efetiva': '{:.2f}%', 'lucro_apos': 'R$ {:,.2f}'
            }))
            with st.expander("Detalhes 2027"):
                for r in res27.values():
                    st.write(f"**{r['regime']}**")
                    for k, v in r['detalhes'].items():
                        st.write(f"- {k}: R$ {v:,.2f}")

            st.subheader("💳 Split Payment CBS")
            c1, c2, c3 = st.columns(3)
            c1.metric("CBS retida/mês", f"R$ {sp['split']:,.2f}")
            c2.metric("Receb. Líquido", f"R$ {sp['receb_liq']:,.2f}")
            c3.metric("Capital de giro", f"R$ {sp['nec_liquida']:,.2f}")

            st.session_state['resultados'] = (res26, res27, sp)

    with tab2:
        if 'resultados' in st.session_state:
            res26, res27, sp = st.session_state['resultados']
            col1, col2 = st.columns(2)
            with col1:
                if st.button("📄 Gerar Relatório PDF"):
                    rel = RelatorioPDF(empresa, res26, res27, sp)
                    pdf_bytes = rel.gerar()
                    st.download_button("⬇️ Baixar PDF", data=pdf_bytes, file_name="relatorio.pdf", mime="application/pdf")
            with col2:
                if st.button("📊 Gerar Slides PowerPoint"):
                    sl = SlidesPPT(empresa, res26, res27, sp)
                    ppt_bytes = sl.gerar()
                    st.download_button("⬇️ Baixar PowerPoint", data=ppt_bytes, file_name="slides.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    st.markdown("---")
    st.caption(f"{ESCRITORIO['nome']} – {ESCRITORIO['contador']} – {ESCRITORIO['crc']} – {ESCRITORIO['endereco']} – Tel: {ESCRITORIO['telefone']}")

if __name__ == "__main__":
    main()